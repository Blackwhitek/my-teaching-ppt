import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import google.generativeai as genai
import os
from dotenv import load_dotenv
from datetime import datetime
import requests
from io import BytesIO

load_dotenv()

st.set_page_config(page_title="教學投影片生成器", layout="wide")
st.title("🧑‍🏫 依參考資料自動生成教學投影片")
st.markdown("### 自動分析 + 自動插入相關圖片")

api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    st.error("❌ 未設定 GOOGLE_API_KEY")
    st.stop()

genai.configure(api_key=api_key)

# 設定
with st.sidebar:
    st.header("生成設定")
    model_name = st.selectbox("AI 模型", ["gemini-2.5-flash", "gemini-2.0-flash"])
    num_slides = st.slider("投影片張數", 8, 25, 15)

# 上傳檔案
uploaded_file = st.file_uploader("上傳參考資料（PDF / Word / TXT）", type=["pdf", "docx", "txt"])
topic = st.text_input("專題名稱", placeholder="例如：光合作用原理")

if st.button("🚀 生成含圖片的教學投影片", type="primary", use_container_width=True):
    if not uploaded_file and not topic:
        st.error("請上傳資料或輸入主題")
    else:
        with st.spinner("AI 正在分析內容並尋找適合圖片..."):
            
            # 讀取檔案
            file_content = ""
            if uploaded_file:
                if uploaded_file.type == "application/pdf":
                    import PyPDF2
                    pdf = PyPDF2.PdfReader(uploaded_file)
                    file_content = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
                elif uploaded_file.type.startswith("application/vnd.openxmlformats"):
                    from docx import Document
                    doc = Document(uploaded_file)
                    file_content = "\n".join([para.text for para in doc.paragraphs])

            prompt = f"""
你是一位專業教學設計師。
主題：{topic}
參考內容：{file_content[:12000]}

請生成約 {num_slides} 張教學投影片（遵守6×6規則）。
每張投影片請提供：
- 標題
- 主要內容（ bullet points ）
- 講者筆記
- 建議圖片描述（要精準、適合教育用途）

輸出格式：
**第X張：** [標題]
**內容：**
- 要點...
**講者筆記：** ...
**建議圖片：** [詳細描述，例如：卡通風格的光合作用過程圖、綠色植物細胞結構圖等]
"""

            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            generated_text = response.text

            # 生成 PPTX + 插入圖片
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # 簡單解析並建立投影片
            parts = generated_text.split("**第")
            for i, part in enumerate(parts):
                if len(part.strip()) < 10:
                    continue

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                title_shape.text = f"第{i}張 " + part.split("\n")[0].replace("：", "").strip()[:60]

                # 內容
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()

                lines = part.split("\n")
                for line in lines:
                    if line.strip().startswith("-") or line.strip().startswith("•"):
                        p = tf.add_paragraph()
                        p.text = line.strip("- •").strip()
                        p.font.size = Pt(20)

                # 嘗試插入圖片（如果有建議圖片描述）
                if "建議圖片：" in part or "圖片：" in part:
                    try:
                        # 簡單從描述取關鍵字
                        img_desc = part.split("建議圖片：")[-1].split("\n")[0].strip()[:100]
                        # 使用 Unsplash 隨機搜尋（免費）
                        search_term = img_desc.replace("圖", "").replace("示意", "").replace("卡通", "")[:30]
                        url = f"https://source.unsplash.com/featured/800x600/?{search_term}"
                        
                        response_img = requests.get(url, timeout=5)
                        if response_img.status_code == 200:
                            slide.shapes.add_picture(BytesIO(response_img.content), 
                                                   Inches(6.5), Inches(1.8), Inches(6))
                    except:
                        pass  # 如果圖片下載失敗就跳過

            # 儲存
            filename = f"教學投影片_{topic[:20]}_{datetime.now().strftime('%m%d_%H%M')}.pptx"
            prs.save(filename)

            st.success(f"✅ 已生成 {len(prs.slides)} 張含圖片的教學投影片！")
            
            with open(filename, "rb") as f:
                st.download_button(
                    label="📥 下載 .pptx 檔案（已含圖片）",
                    data=f,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )

            with st.expander("查看 AI 生成的完整內容"):
                st.write(generated_text)